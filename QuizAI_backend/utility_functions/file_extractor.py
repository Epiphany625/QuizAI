import os
import sys
import json
from PyPDF2 import PdfReader
from pptx import Presentation
from PIL import Image
import pytesseract

def extract_text_from_pdf(filepath):
    """Extracts text from a PDF file."""
    text_content = []
    with open(filepath, 'rb') as f:
        reader = PdfReader(f)
        for page in reader.pages:
            page_text = page.extract_text() or ""
            text_content.append(page_text)
    return "\n".join(text_content)

def extract_text_from_pptx(filepath):
    """Extracts text from a PPTX file."""
    prs = Presentation(filepath)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

def extract_text_from_image(filepath):
    """Extracts text from an image file using OCR."""
    # Open image
    img = Image.open(filepath)
    # Use tesseract to do OCR on the image
    text = pytesseract.image_to_string(img)
    return text

def extract_text(filepath):
    """Determines file type and extracts text accordingly."""
    _, ext = os.path.splitext(filepath)
    ext = ext.lower()

    if ext == ".pdf":
        return extract_text_from_pdf(filepath)
    elif ext == ".pptx":
        return extract_text_from_pptx(filepath)
    elif ext in [".png", ".jpg", ".jpeg", ".tiff", ".bmp"]:
        return extract_text_from_image(filepath)
    else:
        raise ValueError(f"Unsupported file extension: {ext}")

if __name__ == "__main__":
    if len(sys.argv) != 2:
        print(json.dumps({"error": "Usage: python file_extractor.py <file_path>"}))
        sys.exit(1)
    
    file_path = sys.argv[1]
    if not os.path.exists(file_path):
        print(json.dumps({"error": f"File {file_path} does not exist."}))
        sys.exit(1)

    try:
        extracted_text = extract_text(file_path)
        result = {
            "status": "success",
            "extracted_text": extracted_text
        }
        print(json.dumps(result))  # Output JSON result
    except Exception as e:
        print(json.dumps({"status": "error", "message": str(e)}))
        sys.exit(1)
